# apps/document_qa.py
import os
import argparse
import glob
from typing import Dict, List, Any, Tuple, Optional
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import re

from ai_gateway.client import AIGatewayClient
from vector_db.pgvector_client import PGVectorClient
from agents.rag_agent import query_rag_agent

from dotenv import load_dotenv
load_dotenv()

def extract_text_from_pdf(file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """Extract text from a PDF file."""
    chunks = []
    doc = fitz.open(file_path)
    
    for page_num, page in enumerate(doc):
        text = page.get_text()
        
        # Split into paragraphs and chunks
        paragraphs = text.split('\n\n')
        for i, para in enumerate(paragraphs):
            if len(para.strip()) > 10:  # Skip very short paragraphs
                chunks.append((
                    para.strip(),
                    {
                        "source": file_path,
                        "page": page_num + 1,
                        "paragraph": i + 1
                    }
                ))
    
    return chunks

def extract_text_from_docx(file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """Extract text from a DOCX file."""
    chunks = []
    doc = docx.Document(file_path)
    
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if len(text) > 10:  # Skip very short paragraphs
            chunks.append((
                text,
                {
                    "source": file_path,
                    "paragraph": i + 1
                }
            ))
    
    return chunks

def extract_text_from_pptx(file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """Extract text from a PPTX file."""
    chunks = []
    prs = Presentation(file_path)
    
    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        text = "\n".join(slide_text)
        if len(text.strip()) > 10:  # Skip very short slides
            chunks.append((
                text.strip(),
                {
                    "source": file_path,
                    "slide": slide_num + 1
                }
            ))
    
    return chunks

def extract_text_from_html(file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """Extract text from an HTML file."""
    chunks = []
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    # Extract text from paragraphs
    for i, para in enumerate(soup.find_all('p')):
        text = para.get_text().strip()
        if len(text) > 10:  # Skip very short paragraphs
            chunks.append((
                text,
                {
                    "source": file_path,
                    "paragraph": i + 1
                }
            ))
    
    # Extract text from headings
    for i, heading in enumerate(soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])):
        text = heading.get_text().strip()
        if len(text) > 5:  # Skip very short headings
            chunks.append((
                text,
                {
                    "source": file_path,
                    "heading": i + 1
                }
            ))
    
    return chunks

def extract_text_from_txt(file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """Extract text from a TXT file."""
    chunks = []
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    # Split into paragraphs
    paragraphs = text.split('\n\n')
    for i, para in enumerate(paragraphs):
        if len(para.strip()) > 10:  # Skip very short paragraphs
            chunks.append((
                para.strip(),
                {
                    "source": file_path,
                    "paragraph": i + 1
                }
            ))
    
    return chunks

def process_document(file_path: str) -> List[Tuple[str, Dict[str, Any]]]:
    """Process a document and extract text chunks with metadata."""
    _, ext = os.path.splitext(file_path.lower())
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext in ['.html', '.htm']:
        return extract_text_from_html(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    else:
        print(f"Unsupported file format: {ext}")
        return []

def ingest_documents(directory_path: str, recursive: bool = True) -> int:
    """
    Ingest documents from a directory into the vector database.
    
    Args:
        directory_path: Path to the directory containing documents
        recursive: Whether to recursively search subdirectories
        
    Returns:
        Number of chunks ingested
    """
    # Initialize clients
    gateway_client = AIGatewayClient()
    pgvector_client = PGVectorClient()
    
    # Find all documents
    pattern = os.path.join(directory_path, '**' if recursive else '*')
    file_paths = []
    
    for ext in ['.pdf', '.docx', '.pptx', '.html', '.htm', '.txt']:
        file_paths.extend(glob.glob(pattern + ext, recursive=recursive))
    
    if not file_paths:
        print(f"No supported documents found in {directory_path}")
        return 0
    
    print(f"Found {len(file_paths)} documents to process")
    
    # Process documents and collect chunks
    all_chunks = []
    all_metadata = []
    
    for file_path in file_paths:
        print(f"Processing {file_path}")
        chunks = process_document(file_path)
        
        texts = [chunk[0] for chunk in chunks]
        metadata = [chunk[1] for chunk in chunks]
        
        all_chunks.extend(texts)
        all_metadata.extend(metadata)
    
    # Generate embeddings in batches of 100
    batch_size = 100
    total_chunks = len(all_chunks)
    total_inserted = 0
    
    for i in range(0, total_chunks, batch_size):
        batch_end = min(i + batch_size, total_chunks)
        batch_chunks = all_chunks[i:batch_end]
        batch_metadata = all_metadata[i:batch_end]
        
        print(f"Generating embeddings for chunks {i+1}-{batch_end} of {total_chunks}")
        batch_embeddings = gateway_client.generate_embeddings(texts=batch_chunks)
        
        print(f"Inserting chunks {i+1}-{batch_end} into vector database")
        doc_ids = pgvector_client.insert_embeddings(
            contents=batch_chunks,
            embeddings=batch_embeddings,
            metadata=batch_metadata
        )
        
        total_inserted += len(doc_ids)
    
    pgvector_client.close()
    print(f"Successfully ingested {total_inserted} document chunks")
    
    return total_inserted

def ask_question(question: str) -> str:
    """
    Ask a question about the ingested documents.
    
    Args:
        question: The question to ask about the documents
        
    Returns:
        The answer to the question
    """
    return query_rag_agent(question)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Document Q&A System")
    subparsers = parser.add_subparsers(dest="command", help="Command to run")
    
    # Ingest command
    ingest_parser = subparsers.add_parser("ingest", help="Ingest documents")
    ingest_parser.add_argument("--directory", required=True, help="Directory containing documents")
    ingest_parser.add_argument("--recursive", action="store_true", help="Recursively search subdirectories")
    
    # Query command
    query_parser = subparsers.add_parser("query", help="Query the system")
    query_parser.add_argument("--question", required=True, help="Question to ask")
    
    args = parser.parse_args()
    
    if args.command == "ingest":
        ingest_documents(args.directory, args.recursive)
    elif args.command == "query":
        answer = ask_question(args.question)
        print(f"Q: {args.question}")
        print(f"A: {answer}")
    else:
        parser.print_help()